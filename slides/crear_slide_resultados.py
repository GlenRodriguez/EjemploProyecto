from pptx import Presentation
from pptx.util import Inches, Pt
import os

def main():
    metrics_path = "logs/metrics_baseline.txt"
    if not os.path.exists(metrics_path):
        raise FileNotFoundError(f"No se encontró {metrics_path}. Ejecuta primero el baseline.")

    with open(metrics_path, "r", encoding="utf-8") as f:
        metrics_text = f.read()

    prs = Presentation()
    # Slide título
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Resultados Baseline (Dummy + kNN)"
    slide.placeholders[1].text = "Lectura desde logs/metrics_baseline.txt"

    # Slide contenido (métricas)
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Reporte de métricas"
    text_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.text = metrics_text
    for p in tf.paragraphs:
        p.font.size = Pt(12)

    os.makedirs("slides", exist_ok=True)
    out = "slides/Resultados_Baseline_Generado.pptx"
    prs.save(out)
    print(f"Presentación creada: {out}")

if __name__ == "__main__":
    main()
